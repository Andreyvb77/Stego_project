#!/usr/bin/env python3
"""Create PowerPoint presentation for StegoProject VKR pre-defense"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Create presentation
prs = Presentation()
prs.slide_width = Cm(33.87)
prs.slide_height = Cm(19.05)

# Colors
PRIMARY = RGBColor(41, 98, 255)
SECONDARY = RGBColor(52, 152, 219)
ACCENT = RGBColor(231, 76, 60)
SUCCESS = RGBColor(39, 174, 96)
TEXT = RGBColor(44, 62, 80)
LIGHT_BG = RGBColor(245, 247, 250)

print("Creating slides...")

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Система обнаружения стеганографически скрытой информации\nс помощью нейронных сетей"
title.text_frame.paragraphs[0].font.size = Pt(32)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = PRIMARY

subtitle = slide.placeholders[1]
subtitle.text = "Предзащита ВКР\nСтудент группы ФКИБ 22-04\nНаучный руководитель: [ФИО]"
subtitle.text_frame.paragraphs[0].font.size = Pt(18)
print("Slide 1: Title - done")

# Slide 2: Актуальность и цель
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Актуальность и цель работы"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = PRIMARY

body = slide.placeholders[1].text_frame
body.clear()
content = [
    ("Актуальность:", True, 0),
    ("Рост объемов мультимедийного трафика в интернете", False, 1),
    ("Изображения — популярный контейнер для скрытых каналов", False, 1),
    ("Традиционные методы неэффективны против современных алгоритмов", False, 1),
    ("Необходим переход к методам глубокого обучения", False, 1),
    ("", False, 0),
    ("Цель работы:", True, 0),
    ("Разработка и сравнительный анализ CNN архитектур (SRNet, GSR-Net, Yedroudj-Net, Zhu-Net)", False, 1),
    ("Обнаружение следов LSB-стеганографии в растровых изображениях", False, 1),
    ("Выбор оптимальной модели для практического применения", False, 1),
]
for i, (text, bold, level) in enumerate(content):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT
    p.level = level
    if bold:
        p.font.bold = True
print("Slide 2: Актуальность - done")

# Slide 3: Этапы подготовки
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Этапы подготовки и выполнения работы"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = PRIMARY

body = slide.placeholders[1].text_frame
body.clear()
content = [
    ("1. Обзор литературы и методов:", True, 0),
    ("Изучены методы стеганоанализа и архитектуры нейросетей", False, 1),
    ("Проанализированы: SRNet, Yedroudj-Net, Zhu-Net, GSR-Net", False, 1),
    ("", False, 0),
    ("2. Подготовка датасетов:", True, 0),
    ("Cats vs Dogs (1000 изображений) → grayscale", False, 1),
    ("Синтетические градиентные изображения (256×256)", False, 1),
    ("Внедрение скрытого сообщения методом LSB-replacement", False, 1),
    ("", False, 0),
    ("3. Реализация моделей на PyTorch:", True, 0),
    ("Адаптация архитектур для бинарной классификации", False, 1),
    ("Реализация слоев HPF/SRM фильтров", False, 1),
    ("Добавление механизмов внимания (SE-Block, CBAM)", False, 1),
]
for i, (text, bold, level) in enumerate(content):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT
    p.level = level
    if bold:
        p.font.bold = True
print("Slide 3: Этапы - done")

# Slide 4: Метод LSB
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Метод LSB-стеганографии"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

# Original pixel
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(3), Cm(4), Cm(2))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BG
shape.line.color.rgb = PRIMARY
tf = shape.text_frame
tf.text = "Original Pixel\n10110110 (182)"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.size = Pt(14)

# Arrow
slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(7.5), Cm(3.75), Cm(2), Cm(0.5))

# Secret bit
secret_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(8), Cm(4.5), Cm(1), Cm(0.8))
secret_box.fill.solid()
secret_box.fill.fore_color.rgb = ACCENT
secret_tf = secret_box.text_frame
secret_tf.text = "Bit: 1"
secret_tf.paragraphs[0].font.size = Pt(10)
secret_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Modified pixel
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(11), Cm(3), Cm(4), Cm(2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 240, 230)
shape.line.color.rgb = ACCENT
tf = shape.text_frame
tf.text = "Modified Pixel\n10110111 (183)\nLSB changed!"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.size = Pt(12)

# Explanation
exp = slide.shapes.add_textbox(Cm(3), Cm(6), Cm(12), Cm(2))
tf = exp.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "LSB Method:\n• Replace Least Significant Bit with secret data\n• Minimal visual change (+/-1 intensity value)\n• Detectable by neural networks"
p.font.size = Pt(14)
p.font.color.rgb = TEXT
print("Slide 4: LSB - done")

# Slide 5: Архитектура SRNet
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Архитектура SRNet (Steganalysis ResNet)"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

# Diagram blocks
y = Cm(3)
blocks = [
    (Cm(1.5), "Input\n(1x256x256)", LIGHT_BG, PRIMARY),
    (Cm(5), "HPF Layer\n(5 SRM filters)", RGBColor(230, 240, 255), SECONDARY),
    (Cm(9), "Conv1\n(5->64 ch)", RGBColor(230, 240, 255), SECONDARY),
]
for x, text, fill, line in blocks:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Cm(3), Cm(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    tf = shape.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(11)

# Residual blocks
y2 = Cm(5.5)
for i in range(2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.5 + i*5), y2, Cm(4.5), Cm(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 235, 230)
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.text = "ResBlock x3\nstride=2"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(11)

y3 = Cm(8)
for i in range(2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.5 + i*5), y3, Cm(4.5), Cm(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 235, 230)
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.text = "ResBlock x3\nstride=2"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(11)

# Output
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(10.5), Cm(10.5), Cm(4), Cm(1.5))
shape.fill.solid()
shape.fill.fore_color.rgb = SUCCESS
shape.line.color.rgb = SUCCESS
tf = shape.text_frame
tf.text = "Global AvgPool\n+ FC (512->2)"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.size = Pt(12)

# Technical details
details = slide.shapes.add_textbox(Cm(16), Cm(3), Cm(17), Cm(9))
tf = details.text_frame
tf.word_wrap = True
tech = [
    "Техническая реализация:",
    "* HPF Layer: 5 фиксированных SRM-фильтров",
    "* Conv1: 5->64 канала, kernel 3x3",
    "* 4 группы остаточных блоков: 64->64->128->256->512",
    "* ResidualBlock: abs() + tanh() активации",
    "* Skip-connections для борьбы с затуханием градиентов",
    "* Global Average Pooling + FC (512->2)",
    "* Обучаемых параметров: ~2.5 млн",
]
for i, text in enumerate(tech):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(12)
    if i == 0:
        p.font.bold = True
print("Slide 5: SRNet - done")

# Slide 6: Архитектура GSR-Net
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Архитектура GSR-Net (с механизмом внимания)"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

# Diagram
y = Cm(2.5)
blocks = [
    (Cm(1), "Input\n(1x256x256)", LIGHT_BG, PRIMARY),
    (Cm(4), "HPF\n(5 filters)", RGBColor(230, 240, 255), SECONDARY),
    (Cm(7), "Conv1\n(5->64)", RGBColor(230, 240, 255), SECONDARY),
]
for x, text, fill, line in blocks:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Cm(2.5), Cm(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    tf = shape.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(10)

# Attention blocks
y2 = Cm(4.5)
channels = ["64->64", "64->128", "128->256", "256->512"]
for i, ch in enumerate(channels):
    x = Cm(0.5 + i*4.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y2, Cm(4), Cm(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(230, 250, 235)
    shape.line.color.rgb = SUCCESS
    tf = shape.text_frame
    tf.text = f"SE-Block\n+ ResBlock x3\n{ch}"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(9)
    
    # SE indicator
    se = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x+1.5, y2+1.7, Cm(1), Cm(1))
    se.fill.solid()
    se.fill.fore_color.rgb = ACCENT
    se_tf = se.text_frame
    se_tf.text = "SE"
    se_tf.paragraphs[0].font.size = Pt(9)
    se_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Output
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(13.5), Cm(7.5), Cm(3.5), Cm(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = SUCCESS
shape.line.color.rgb = SUCCESS
tf = shape.text_frame
tf.text = "AvgPool + FC\n(512->2)"
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.paragraphs[0].font.size = Pt(11)

# Details
details = slide.shapes.add_textbox(Cm(1), Cm(9.5), Cm(32), Cm(8))
tf = details.text_frame
tf.word_wrap = True
tech = [
    "Ключевые особенности GSR-Net:",
    "* SE-Block (Squeeze-and-Excitation): канальная рекалибровка",
    "  - Squeeze: глобальное усреднение по пространству",
    "  - Excitation: обучение весов важности каналов",
    "  - Scale: умножение признаков на веса",
    "* AttentionResidualBlock: SE-Block встроен в остаточный блок",
    "* Фокусировка на информативных признаках стеганографических искажений",
    "* Обучаемых параметров: ~3.1 млн",
]
for i, text in enumerate(tech):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(12)
    if i == 0:
        p.font.bold = True
print("Slide 6: GSR-Net - done")

# Slide 7: Yedroudj-Net и Zhu-Net
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Архитектуры Yedroudj-Net и Zhu-Net"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = PRIMARY

body = slide.placeholders[1].text_frame
body.clear()
content = [
    ("Yedroudj-Net (базовая архитектура):", True, 0),
    ("5 сверточных блоков с Tanh активацией и AvgPooling", False, 1),
    ("Последовательность: 1->64->16->16->16->16 каналов", False, 1),
    ("3 полносвязных слоя: 1024->512->256->2", False, 1),
    ("Dropout 0.5 для регуляризации", False, 1),
    ("~1.8 млн обучаемых параметров", False, 1),
    ("", False, 0),
    ("Zhu-Net (State-of-the-Art):", True, 0),
    ("30 SRM-фильтров для извлечения высокочастотных остатков", False, 1),
    ("AttentionBlock в каждом остаточном блоке", False, 1),
    ("4 группы residual блоков: 64->128->256->512", False, 1),
    ("Глобальный средний пулинг + классификатор", False, 1),
    ("~4.2 млн обучаемых параметров", False, 1),
]
for i, (text, bold, level) in enumerate(content):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = text
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT
    p.level = level
    if bold:
        p.font.bold = True
print("Slide 7: Yedroudj/Zhu - done")

# Slide 8: Результаты экспериментов
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(32), Cm(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Результаты экспериментов (тестовый набор)"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

# Chart
models = ['SRNet', 'GSR-Net', 'Yedroudj-Net', 'Zhu-Net']
accuracy = [0.7233, 0.7400, 0.6033, 0.6867]
auc = [0.7983, 0.8119, 0.6629, 0.8298]

chart_data = CategoryChartData()
chart_data.categories = models
chart_data.add_series('Accuracy', accuracy)
chart_data.add_series('ROC-AUC', auc)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Cm(2), Cm(3), Cm(28), Cm(12), chart_data
).chart
chart.has_legend = True
print("Slide 8: Results chart - done")

# Slide 9: Сравнительный анализ
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Сравнительный анализ эффективности"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = PRIMARY

body = slide.placeholders[1].text_frame
body.clear()
content = [
    ("Таблица результатов:", True, 0),
    ("", False, 0),
    ("Модель          | Accuracy | ROC-AUC | Параметры", True, 0),
    ("----------------------------------------------", False, 0),
    ("GSR-Net         |  74.00%  |  0.8119  | ~3.1 млн * Лучшая accuracy", False, 0),
    ("SRNet           |  72.33%  |  0.7983  | ~2.5 млн * Баланс", False, 0),
    ("Zhu-Net         |  68.67%  |  0.8298  | ~4.2 млн * Лучший AUC", False, 0),
    ("Yedroudj-Net    |  60.33%  |  0.6629  | ~1.8 млн", False, 0),
    ("", False, 0),
    ("Выводы:", True, 0),
    ("- GSR-Net показал наилучшую точность детектирования (74%)", False, 1),
    ("- SE-Block улучшает качество на ~2% по сравнению с SRNet", False, 1),
    ("- Zhu-Net имеет лучший ROC-AUC (0.83), но требует больше ресурсов", False, 1),
    ("- Все модели превосходят случайное угадывание (50%)", False, 1),
]
for i, (text, bold, level) in enumerate(content):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT
    p.level = level
    if bold:
        p.font.bold = True
print("Slide 9: Comparison - done")

# Slide 10: План дальнейшей работы
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "План дальнейшей работы"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = PRIMARY

body = slide.placeholders[1].text_frame
body.clear()
content = [
    ("1. Расширение экспериментальной базы:", True, 0),
    ("Тестирование на больших датасетах (BOSSbase, BOWS-2)", False, 1),
    ("Исследование устойчивости к другим методам стеганографии", False, 1),
    ("Оценка работы при разном уровне внедрения", False, 1),
    ("", False, 0),
    ("2. Оптимизация моделей:", True, 0),
    ("Pruning и квантование для уменьшения размера", False, 1),
    ("Исследование transfer learning", False, 1),
    ("", False, 0),
    ("3. Разработка программного модуля:", True, 0),
    ("Создание веб-интерфейса для анализа изображений", False, 1),
    ("Интеграция лучшей модели (GSR-Net) в систему", False, 1),
    ("Подготовка документации", False, 1),
    ("", False, 0),
    ("4. Написание ВКР:", True, 0),
    ("Оформление теоретической и практической частей", False, 1),
    ("Подготовка демонстрационных материалов", False, 1),
]
for i, (text, bold, level) in enumerate(content):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT
    p.level = level
    if bold:
        p.font.bold = True
print("Slide 10: Future work - done")

# Slide 11: Заключение
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Заключение"
title.text_frame.paragraphs[0].font.size = Pt(28)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = PRIMARY

body = slide.placeholders[1].text_frame
body.clear()
content = [
    ("Выполненные задачи:", True, 0),
    ("- Проведен обзор методов стеганоанализа и нейросетевых архитектур", False, 1),
    ("- Подготовлены датасеты (Cats vs Dogs grayscale, синтетические)", False, 1),
    ("- Реализованы 4 архитектуры нейронных сетей на PyTorch", False, 1),
    ("- Внедрен метод LSB-стеганографии для генерации данных", False, 1),
    ("- Проведены эксперименты и сравнительный анализ", False, 1),
    ("", False, 0),
    ("Научная новизна:", True, 0),
    ("Адаптация специализированных архитектур стегоанализа", False, 1),
    ("для черно-белых изображений и оценка их эффективности", False, 1),
    ("", False, 0),
    ("Практическая значимость:", True, 0),
    ("Прототип системы обнаружения скрытой информации", False, 1),
    ("Рекомендации по выбору оптимальной архитектуры", False, 1),
]
for i, (text, bold, level) in enumerate(content):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT
    p.level = level
    if bold:
        p.font.bold = True
print("Slide 11: Conclusion - done")

# Slide 12: Спасибо
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.87), Cm(19.05))
bg.fill.solid()
bg.fill.fore_color.rgb = LIGHT_BG
bg.line.fill.background()

text_box = slide.shapes.add_textbox(Cm(5), Cm(7), Cm(24), Cm(5))
tf = text_box.text_frame
p = tf.paragraphs[0]
p.text = "Спасибо за внимание!"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "\nГотов ответить на ваши вопросы"
p2.font.size = Pt(20)
p2.font.color.rgb = TEXT
p2.alignment = PP_ALIGN.CENTER
print("Slide 12: Thank you - done")

# Save
output_path = '/workspace/StegoProject_Presentation.pptx'
prs.save(output_path)
print(f"\n=== Presentation saved to: {output_path} ===")
print(f"Total slides: {len(prs.slides)}")
